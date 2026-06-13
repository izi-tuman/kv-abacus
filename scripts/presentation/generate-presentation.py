"""
Создание продающей PPTX презентации для B2B — система управления бронированием.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Цветовая палитра ──
ACCENT = RGBColor(0x1B, 0x7A, 0x5E)       # Основной акцент (зелёный)
ACCENT_LIGHT = RGBColor(0xE8, 0xF5, 0xF0)  # Светлый фон
ACCENT_DARK = RGBColor(0x0F, 0x4C, 0x3A)   # Тёмный акцент
DARK = RGBColor(0x1A, 0x1A, 0x2E)          # Тёмный текст
MUTED = RGBColor(0x6B, 0x72, 0x80)         # Приглушённый текст
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF9, 0xFA, 0xFB)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)
DANGER = RGBColor(0xDC, 0x26, 0x26)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

SCRIPT_DIR = os.path.dirname(__file__)
SCREENSHOTS_DIR = os.path.join(SCRIPT_DIR, "screenshots")


def add_background(slide, color):
    """Заливка фона слайда."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """Добавить фигуру."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    """Скруглённый прямоугольник."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Добавить текстовое поле."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=14, color=DARK, bold=False, space_before=Pt(4), space_after=Pt(2), alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Добавить параграф в text_frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.space_after = space_after
    p.alignment = alignment
    return p


def add_icon_circle(slide, left, top, size, fill_color, icon_text, font_size=28):
    """Круг с текстом-иконкой."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(font_size)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return shape


def add_screenshot_placeholder(slide, left, top, width, height, label, filename=None):
    """Плейсхолдер для скриншота."""
    shape = add_rounded_rect(slide, left, top, width, height, fill_color=LIGHT_GRAY, line_color=BORDER)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "📱"
    p.font.size = Pt(36)
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(20)
    add_paragraph(tf, label, font_size=13, color=MUTED, alignment=PP_ALIGN.CENTER)
    if filename:
        add_paragraph(tf, f"Файл: {filename}", font_size=10, color=RGBColor(0x9C, 0xA3, 0xAF), alignment=PP_ALIGN.CENTER)

    # Попытка вставить реальный скриншот
    if filename:
        screenshot_path = os.path.join(SCREENSHOTS_DIR, filename)
        if os.path.exists(screenshot_path):
            try:
                # Удаляем плейсхолдер
                sp = shape._element
                sp.getparent().remove(sp)
                slide.shapes.add_picture(screenshot_path, left, top, width, height)
            except Exception:
                pass
    return shape


def add_accent_bar(slide, left, top, width, height):
    """Акцентная полоска."""
    add_shape(slide, left, top, width, height, fill_color=ACCENT)


def add_top_accent_line(slide):
    """Тонкая линия сверху."""
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Pt(4), fill_color=ACCENT)


def add_slide_number(slide, num, total=16):
    """Номер слайда."""
    add_text_box(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
                 f"{num}/{total}", font_size=10, color=MUTED, alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════
# СЛАЙД 1: Титульный
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_background(slide, DARK)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Pt(6))

# Декоративные элементы
add_shape(slide, Inches(9), Inches(0.5), Inches(4), Inches(4), fill_color=ACCENT)
from pptx.util import Emu
shape_el = slide.shapes[-1]._element
from lxml import etree
# Скругляем через rotation
shape_el.set("rotation", "45")

add_text_box(slide, Inches(1), Inches(1.5), Inches(8), Inches(1.2),
             "СИСТЕМА УПРАВЛЕНИЯ", font_size=20, color=ACCENT, bold=True, font_name="Calibri Light")
add_text_box(slide, Inches(1), Inches(2.3), Inches(10), Inches(1.5),
             "БРОНИРОВАНИЕМ ДОМОВ", font_size=44, color=WHITE, bold=True, font_name="Calibri Light")

add_text_box(slide, Inches(1), Inches(4.0), Inches(8), Inches(0.8),
             "Автоматизация вашего гостиничного бизнеса", font_size=22, color=RGBColor(0xA0, 0xB0, 0xC0), font_name="Calibri Light")

add_text_box(slide, Inches(1), Inches(5.2), Inches(8), Inches(0.6),
             "Единая платформа для управления объектами, клиентами и бронированиями",
             font_size=16, color=MUTED, font_name="Calibri Light")

# Линия-разделитель
add_shape(slide, Inches(1), Inches(6.0), Inches(3), Pt(2), fill_color=ACCENT)

add_slide_number(slide, 1)

# ═══════════════════════════════════════════════
# СЛАЙД 2: О системе
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_top_accent_line(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
             "О СИСТЕМЕ", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(1.0),
             "Всё для управления арендой домов в одном месте", font_size=32, color=DARK, bold=True)

add_text_box(slide, Inches(0.8), Inches(2.3), Inches(7), Inches(1.5),
             "Система создана для владельцев баз отдыха, глэмпингов и гостиничного бизнеса. "
             "Она объединяет управление объектами, клиентской базой, бронированиями и аналитикой "
             "в едином интуитивном интерфейсе, доступном с любого устройства.",
             font_size=16, color=MUTED)

# 4 карточки-преимущества
cards = [
    ("🏠", "Управление\nобъектами", "Дома, вместимость,\nстоимость"),
    ("📅", "Календарь\nбронирований", "Неделя, день,\nконфликты"),
    ("👥", "Клиентская\nбаза", "Контакты, история,\nзаметки"),
    ("📊", "Отчёты\nи аналитика", "Финансы, загрузка,\nменеджеры"),
]
for i, (icon, title, desc) in enumerate(cards):
    x = Inches(0.8) + Inches(3.0) * i
    card = add_rounded_rect(slide, x, Inches(4.0), Inches(2.7), Inches(2.8), fill_color=LIGHT_GRAY, line_color=BORDER)
    add_text_box(slide, x + Inches(0.2), Inches(4.2), Inches(2.3), Inches(0.6),
                 icon, font_size=32, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(4.9), Inches(2.3), Inches(0.8),
                 title, font_size=16, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(5.7), Inches(2.3), Inches(0.6),
                 desc, font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 2)

# ═══════════════════════════════════════════════
# СЛАЙД 3: Управление домами
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_top_accent_line(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
             "01", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
             "Управление домами", font_size=32, color=DARK, bold=True)

# Описание
desc_box = add_text_box(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.5),
    "• Добавление объектов с указанием вместимости и стоимости\n"
    "• Описание и фотографии каждого дома\n"
    "• Активация/деактивация объектов\n"
    "• Быстрый просмотр списка с ценами\n"
    "• Фильтрация по статусу",
    font_size=15, color=MUTED)

# Скриншот
add_screenshot_placeholder(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2),
    "Скриншот: Список домов с ценами и вместимостью", "houses.png")

add_slide_number(slide, 3)

# ═══════════════════════════════════════════════
# СЛАЙД 4: Календарь бронирований — Недельный вид
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_top_accent_line(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
             "02", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
             "Календарь бронирований — Недельный вид", font_size=32, color=DARK, bold=True)

desc_box = add_text_box(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.5),
    "• Обзор всех 7 дней недели на одном экране\n"
    "• Визуальное отображение занятых и свободных дат\n"
    "• Карточки бронирований с ключевой информацией\n"
    "• Быстрое создание брони кликом на свободный слот\n"
    "• Навигация между неделями и месяцами",
    font_size=15, color=MUTED)

add_screenshot_placeholder(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2),
    "Скриншот: Недельный вид календаря", "calendar-week.png")

add_slide_number(slide, 4)

# ═══════════════════════════════════════════════
# СЛАЙД 5: Календарь — Дневной вид + Модалка
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_top_accent_line(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
             "02", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
             "Дневной вид и создание бронирования", font_size=32, color=DARK, bold=True)

desc_box = add_text_box(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.5),
    "• Детальный просмотр одного дня\n"
    "• Модальное окно с формой бронирования\n"
    "• Выбор клиента и дома через поиск\n"
    "• Выбор дат через календарь\n"
    "• Добавление доп. услуг и предоплаты",
    font_size=15, color=MUTED)

add_screenshot_placeholder(slide, Inches(6.8), Inches(1.5), Inches(2.7), Inches(5.2),
    "Дневной вид", "calendar-day.png")
add_screenshot_placeholder(slide, Inches(9.8), Inches(1.5), Inches(2.7), Inches(5.2),
    "Модалка бронирования", "booking-modal.png")

add_slide_number(slide, 5)

# ═══════════════════════════════════════════════
# СЛАЙД 6: Клиентская база
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_top_accent_line(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
             "03", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
             "Клиентская база", font_size=32, color=DARK, bold=True)

desc_box = add_text_box(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(3.0),
    "• База клиентов с контактами (имя, телефон, email)\n"
    "• История бронирований каждого клиента\n"
    "• История проката снаряжения\n"
    "• Примечания и заметки\n"
    "• Поиск по имени или телефону\n"
    "• Быстрый вызов по номеру телефона\n"
    "• Чёрный список",
    font_size=15, color=MUTED)

add_screenshot_placeholder(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2),
    "Скриншот: Клиентская база", "clients.png")

add_slide_number(slide, 6)

# ═══════════════════════════════════════════════
# СЛАЙД 7: Услуги
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_top_accent_line(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
             "04", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
             "Дополнительные услуги", font_size=32, color=DARK, bold=True)

desc_box = add_text_box(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.5),
    "• Управление доп. услугами (сауна, чан, баня и т.д.)\n"
    "• Привязка услуг к бронированию\n"
    "• Активация/деактивация услуг\n"
    "• Стоимость каждой услуги\n"
    "• Быстрый поиск при добавлении в бронь",
    font_size=15, color=MUTED)

add_screenshot_placeholder(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2),
    "Скриншот: Управление услугами", "services.png")

add_slide_number(slide, 7)

# ═══════════════════════════════════════════════
# СЛАЙД 8: Снаряжение
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_top_accent_line(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
             "05", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
             "Учёт снаряжения", font_size=32, color=DARK, bold=True)

desc_box = add_text_box(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.5),
    "• Каталог инвентаря (катамаран, лыжи, снегоход)\n"
    "• Управление арендой снаряжения\n"
    "• Учёт количества единиц\n"
    "• История прокатов по каждому клиенту\n"
    "• Стоимость проката",
    font_size=15, color=MUTED)

add_screenshot_placeholder(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2),
    "Скриншот: Управление снаряжением", "equipment.png")

add_slide_number(slide, 8)

# ═══════════════════════════════════════════════
# СЛАЙД 9: Прокат снаряжения
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_top_accent_line(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
             "05", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
             "Прокат снаряжения", font_size=32, color=DARK, bold=True)

desc_box = add_text_box(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.5),
    "• Создание записей о прокате\n"
    "• Выбор клиента из базы\n"
    "• Выбор дат начала и окончания\n"
    "• Добавление нескольких единиц снаряжения\n"
    "• Фильтрация и поиск по прокатам",
    font_size=15, color=MUTED)

add_screenshot_placeholder(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2),
    "Скриншот: Страница проката", "rental.png")

add_slide_number(slide, 9)

# ═══════════════════════════════════════════════
# СЛАЙД 10: История бронирований
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_top_accent_line(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
             "06", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
             "История бронирований", font_size=32, color=DARK, bold=True)

desc_box = add_text_box(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.5),
    "• Архив всех бронирований\n"
    "• Фильтрация по датам (от/до)\n"
    "• Фильтрация по дому\n"
    "• Фильтрация по менеджеру\n"
    "• Фильтрация по статусу (активно/завершено)\n"
    "• Поиск по клиенту",
    font_size=15, color=MUTED)

add_screenshot_placeholder(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2),
    "Скриншот: История бронирований", "history.png")

add_slide_number(slide, 10)

# ═══════════════════════════════════════════════
# СЛАЙД 11: Отчёты — Финансовый
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_top_accent_line(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
             "07", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
             "Финансовый отчёт", font_size=32, color=DARK, bold=True)

desc_box = add_text_box(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.5),
    "• Выручка по дням\n"
    "• Количество бронирований\n"
    "• Средний чек\n"
    "• Настраиваемый период\n"
    "• Детализация по каждому бронированию",
    font_size=15, color=MUTED)

add_screenshot_placeholder(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2),
    "Скриншот: Финансовый отчёт", "reports-finance.png")

add_slide_number(slide, 11)

# ═══════════════════════════════════════════════
# СЛАЙД 12: Отчёты — По менеджерам и домам
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_top_accent_line(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
             "07", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
             "Отчёты по менеджерам и домам", font_size=32, color=DARK, bold=True)

# Левая колонка — менеджеры
add_text_box(slide, Inches(0.8), Inches(2.0), Inches(3), Inches(0.5),
             "По менеджерам", font_size=18, color=DARK, bold=True)
desc_box = add_text_box(slide, Inches(0.8), Inches(2.6), Inches(5), Inches(2.0),
    "• Количество бронирований каждого менеджера\n"
    "• Выручка по каждому менеджеру\n"
    "• Рейтинг эффективности\n"
    "• Настраиваемый период",
    font_size=14, color=MUTED)

# Правая колонка — дома
add_text_box(slide, Inches(7), Inches(2.0), Inches(3), Inches(0.5),
             "По домам", font_size=18, color=DARK, bold=True)
desc_box = add_text_box(slide, Inches(7), Inches(2.6), Inches(5), Inches(2.0),
    "• Загрузка каждого дома\n"
    "• Количество бронирований\n"
    "• Выручка по каждому дому\n"
    "• Процент загрузки",
    font_size=14, color=MUTED)

add_screenshot_placeholder(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.5),
    "Отчёт по менеджерам", "reports-finance.png")
add_screenshot_placeholder(slide, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.5),
    "Отчёт по домам", "reports-finance.png")

add_slide_number(slide, 12)

# ═══════════════════════════════════════════════
# СЛАЙД 13: Преимущества для бизнеса (1)
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Pt(6))

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
             "ПОЧЕМУ ЭТО ВАЖНО", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
             "Преимущества для бизнеса", font_size=36, color=WHITE, bold=True)

benefits = [
    ("⏱", "Экономия времени", "Автоматизация рутинных операций — бронирование за 30 секунд вместо 5 минут"),
    ("📋", "Централизованная база", "Все данные в одном месте: объекты, клиенты, брони, финансы"),
    ("📊", "Контроль в реальном времени", "Видьте загрузку объектов и конфликты бронирований мгновенно"),
    ("📈", "Аналитика для решений", "Финансовые отчёты, загрузка домов, эффективность менеджеров"),
    ("🔐", "Ролевая система", "Разграничение доступа — каждый сотрудник видит только своё"),
]

for i, (icon, title, desc) in enumerate(benefits):
    y = Inches(2.3) + Inches(0.95) * i
    add_icon_circle(slide, Inches(1.0), y, Inches(0.5), ACCENT, icon, font_size=20)
    add_text_box(slide, Inches(1.8), y - Inches(0.05), Inches(4), Inches(0.4),
                 title, font_size=18, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.8), y + Inches(0.35), Inches(10), Inches(0.4),
                 desc, font_size=13, color=RGBColor(0xA0, 0xB0, 0xC0))

add_slide_number(slide, 13)

# ═══════════════════════════════════════════════
# СЛАЙД 14: Преимущества — мобильность
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_top_accent_line(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
             "МОБИЛЬНОСТЬ", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
             "Работайте откуда угодно", font_size=36, color=DARK, bold=True)

add_text_box(slide, Inches(0.8), Inches(2.2), Inches(6), Inches(1.5),
             "Система оптимизирована для мобильных устройств. "
             "Управляйте бронированиями, просматривайте отчёты и работайте с клиентской базой "
             "прямо со смартфона — в любой точке, где есть интернет.",
             font_size=16, color=MUTED)

# 3 карточки
mobile_cards = [
    ("📱", "Мобильная\nверсия", "Адаптивный дизайн\nдля смартфонов"),
    ("🔄", "Pull-to-\nrefresh", "Обновление данных\nсвайпом вниз"),
    ("➕", "Быстрое\nсоздание", "Новая бронь\nв 2 клика"),
]
for i, (icon, title, desc) in enumerate(mobile_cards):
    x = Inches(0.8) + Inches(3.9) * i
    add_rounded_rect(slide, x, Inches(3.8), Inches(3.5), Inches(3.0), fill_color=LIGHT_GRAY, line_color=BORDER)
    add_text_box(slide, x + Inches(0.2), Inches(4.0), Inches(3.1), Inches(0.6),
                 icon, font_size=36, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(4.8), Inches(3.1), Inches(0.8),
                 title, font_size=18, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(5.6), Inches(3.1), Inches(0.6),
                 desc, font_size=13, color=MUTED, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 14)

# ═══════════════════════════════════════════════
# СЛАЙД 15: Для кого подходит
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, ACCENT_LIGHT)
add_top_accent_line(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
             "ДЛЯ КОГО", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
             "Система подходит для вашего бизнеса", font_size=36, color=DARK, bold=True)

targets = [
    ("🏕", "Базы отдыха", "Управление домиками, учёт загрузки,\nбронирование дополнительных услуг"),
    ("⛺", "Глэмпинги", "Контроль занятости, клиентская база,\nотчёты по выручке"),
    ("🏡", "Аренда загородных\nдомов", "Календарь бронирований, история\nклиентов, финансовая аналитика"),
    ("🏨", "Мини-гостиницы", "Управление номерным фондом,\nстатистика, отчёты по менеджерам"),
]

for i, (icon, title, desc) in enumerate(targets):
    row = i // 2
    col = i % 2
    x = Inches(0.8) + Inches(6.0) * col
    y = Inches(2.5) + Inches(2.3) * row

    card = add_rounded_rect(slide, x, y, Inches(5.5), Inches(2.0), fill_color=WHITE, line_color=BORDER)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.2), Inches(1), Inches(0.6),
                 icon, font_size=32)
    add_text_box(slide, x + Inches(1.2), y + Inches(0.2), Inches(4), Inches(0.6),
                 title, font_size=20, color=DARK, bold=True)
    add_text_box(slide, x + Inches(1.2), y + Inches(0.9), Inches(4), Inches(0.8),
                 desc, font_size=13, color=MUTED)

add_slide_number(slide, 15)

# ═══════════════════════════════════════════════
# СЛАЙД 16: Заключение
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Pt(6))

add_text_box(slide, Inches(1), Inches(1.5), Inches(10), Inches(1.2),
             "ГОТОВЫ АВТОМАТИЗИРОВАТЬ\nВАШ БИЗНЕС?", font_size=44, color=WHITE, bold=True)

add_text_box(slide, Inches(1), Inches(3.2), Inches(8), Inches(1.0),
             "Свяжитесь с нами для демонстрации системы и обсуждения условий внедрения",
             font_size=20, color=RGBColor(0xA0, 0xB0, 0xC0))

# CTA кнопка
cta = add_rounded_rect(slide, Inches(1), Inches(4.8), Inches(4), Inches(0.8), fill_color=ACCENT)
tf = cta.text_frame
p = tf.paragraphs[0]
p.text = "ЗАПРОСИТЬ ДЕМО"
p.font.size = Pt(18)
p.font.color.rgb = WHITE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
tf.paragraphs[0].space_before = Pt(0)

add_text_box(slide, Inches(1), Inches(6.0), Inches(6), Inches(0.5),
             "Индивидуальный подход • Быстрое внедрение • Поддержка",
             font_size=14, color=MUTED)

add_slide_number(slide, 16)

# ── Сохранение ──
output_path = os.path.join(SCRIPT_DIR, "Презентация_Система_Бронирования.pptx")
prs.save(output_path)
print(f"\nPresentation saved: {output_path}")
print(f"  Slides: {len(prs.slides)}")
print(f"  Screenshots dir: {SCREENSHOTS_DIR}")
